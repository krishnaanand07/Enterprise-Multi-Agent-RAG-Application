"""
Text extraction logic for various document types.

Supported formats:
  - PDF   (.pdf)  — via pypdf (successor to deprecated PyPDF2)
  - Word  (.docx) — via python-docx
  - PPT   (.pptx) — via python-pptx
  - Excel (.xlsx) — via openpyxl / pandas
  - Plain (.txt, .md, .csv, .json)

All heavy extraction runs in a thread pool (asyncio.to_thread) so the
async event loop is never blocked.
"""
import asyncio
from pathlib import Path
from loguru import logger


class DocumentExtractor:
    """Async-safe extractor that dispatches to format-specific handlers."""

    async def extract_text_async(self, file_path: str) -> tuple[str, dict]:
        """
        Async entry point. Runs the blocking extraction in a thread pool.
        Returns: (extracted_text, metadata_dict)
        """
        return await asyncio.to_thread(self.extract_text, file_path)

    def extract_text(self, file_path: str) -> tuple[str, dict]:
        """
        Synchronous extraction — call via extract_text_async from async code.
        Returns: (extracted_text, metadata_dict)
        """
        ext = Path(file_path).suffix.lower()

        dispatch = {
            ".pdf":  self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
            ".txt":  self._extract_plain,
            ".md":   self._extract_plain,
            ".csv":  self._extract_plain,
            ".json": self._extract_plain,
        }

        handler = dispatch.get(ext)
        if handler is None:
            raise ValueError(
                f"No extractor implemented for '{ext}'. "
                f"Supported: {', '.join(sorted(dispatch.keys()))}"
            )

        return handler(file_path)

    # ------------------------------------------------------------------ #
    #  Plain text                                                          #
    # ------------------------------------------------------------------ #
    def _extract_plain(self, file_path: str) -> tuple[str, dict]:
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read()
                return text, {"num_pages": 1}
            except UnicodeDecodeError:
                continue
        # Last resort: ignore bad chars
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return text, {"num_pages": 1}

    # ------------------------------------------------------------------ #
    #  PDF — uses pypdf (replaces deprecated PyPDF2)                      #
    # ------------------------------------------------------------------ #
    def _extract_pdf(self, file_path: str) -> tuple[str, dict]:
        try:
            from pypdf import PdfReader  # pypdf >= 4.x
        except ImportError:
            # Graceful fallback to PyPDF2 if pypdf not installed yet
            from PyPDF2 import PdfReader  # type: ignore

        text_parts: list[str] = []
        num_pages = 0

        try:
            reader = PdfReader(file_path)
            num_pages = len(reader.pages)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        except Exception as e:
            logger.error(f"PDF extraction error for {file_path}: {e}")

        text = "\n".join(text_parts)

        if not text.strip():
            logger.warning(
                f"No selectable text in {file_path}. "
                "File may be a scanned/image PDF — OCR is required."
            )
            text = "[Image-based PDF: no selectable text found. OCR processing required.]"

        return text, {"num_pages": num_pages}

    # ------------------------------------------------------------------ #
    #  Word (.docx)                                                        #
    # ------------------------------------------------------------------ #
    def _extract_docx(self, file_path: str) -> tuple[str, dict]:
        try:
            from docx import Document
        except ImportError as e:
            raise ImportError("python-docx is required for .docx extraction.") from e

        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        paragraphs.append(row_text)

            text = "\n".join(paragraphs)
            return text, {"num_pages": max(1, len(paragraphs) // 30)}
        except Exception as e:
            logger.error(f"DOCX extraction error for {file_path}: {e}")
            raise

    # ------------------------------------------------------------------ #
    #  PowerPoint (.pptx)                                                  #
    # ------------------------------------------------------------------ #
    def _extract_pptx(self, file_path: str) -> tuple[str, dict]:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx is required for .pptx extraction.") from e

        try:
            prs = Presentation(file_path)
            slide_texts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                    # Extract table cells inside PPT shapes
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            ]
                            if cells:
                                parts.append(" | ".join(cells))

                if parts:
                    slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

            text = "\n\n".join(slide_texts)
            return text, {"num_pages": len(prs.slides)}
        except Exception as e:
            logger.error(f"PPTX extraction error for {file_path}: {e}")
            raise

    # ------------------------------------------------------------------ #
    #  Excel (.xlsx)                                                        #
    # ------------------------------------------------------------------ #
    def _extract_xlsx(self, file_path: str) -> tuple[str, dict]:
        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError("pandas is required for .xlsx extraction.") from e

        try:
            xl = pd.ExcelFile(file_path)
            sheet_texts: list[str] = []

            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                # Convert to readable text representation
                sheet_str = df.to_string(index=False)
                sheet_texts.append(f"[Sheet: {sheet_name}]\n{sheet_str}")

            text = "\n\n".join(sheet_texts)
            return text, {"num_pages": len(xl.sheet_names)}
        except Exception as e:
            logger.error(f"XLSX extraction error for {file_path}: {e}")
            raise


extractor = DocumentExtractor()
